from pptx import Presentation
from pptx.util import Inches, Pt
import os
import pandas as pd
from PIL import Image

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

data = pd.read_csv('Input.csv', encoding = "ISO-8859-1")
attendance = pd.read_csv('Attendance.csv')
prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
ileft = itop = Inches(1)





for index,row in data.iterrows():
    
    try:
        dirname = os.path.dirname(os.path.abspath(__file__))
        os.chdir(os.path.join(dirname, 'headshot'))
    except:
        print("There is an issue with navigating to the headshot folder. Please make sure the folder exists in the main directory.")

    slide = prs.slides.add_slide(blank_slide_layout)
    NetID = str(row['NetID'])
    
    try:
        img_path = NetID + '.jpg'
        im = Image.open(img_path)
    except:
        try:
            img_path = NetID + '.PNG'
            im = Image.open(img_path)
        except:
            try:
                img_path = NetID + '.jpeg'
                im = Image.open(NetID + '.jpeg')
            except:
                try:
                    img_path = NetID + '.JPG'
                    im = Image.open(img_path)
                except:
                    img_path = NetID + '.jp.jpg'
                    im = Image.open(img_path)





    size = im.size

    #scaling the image
    if size[0] > size[1]:
        fac = size[0]/2
        height = size[1] / fac

        iheight = Inches(height)
        iwidth = Inches(2)

    else:
        fac = size[1]/2
        width = size[0] / fac

        iheight = Inches(2)
        iwidth = Inches(width)

    try:
        pic = slide.shapes.add_picture(img_path,ileft,itop,iwidth,iheight)
    except:
        print(NetID + ' headshot image unsupported. skipping...')
        

  
    nameBox = slide.shapes.add_textbox(Inches(1),Inches(0.25),Inches(8),Inches(1))
    nm = nameBox.text_frame
    nm.word_wrap = True
    nm.text = ''
    p = nm.add_paragraph()
    p.text = row['Name'] + '  |  ' + row['NetID'] +'  |  ' + row['Major'] + '  |  ' + row['College'] + ' | ' + str(row['Graduation Year']) 
    p.font.name = 'Garamond'

    
    
    
    try:
        os.chdir(dirname)
        os.chdir(os.path.join(dirname, 'resume_convert'))
    except:
        print("There is an issue with navigating to the Resume folder. Please make sure the folder exists in the main directory.")
    
    try:
        img_path = NetID + '.png'
        im = Image.open(img_path)
    except:
        im = Image.open("blank.png")

    
    size = im.size

    if size[0] > size[1]:
        fac = size[0]/6
        height = size[1] / fac

        iheight = Inches(height)
        iwidth = Inches(6)

    else:
        fac = size[1]/6
        width = size[0] / fac

        iheight = Inches(6)
        iwidth = Inches(width)

    try:
        pic = slide.shapes.add_picture(img_path,Inches(4),itop,iwidth,iheight)
    except:
        print(NetID + ' resume image unsupported. skipping...')

    
    shapes = slide.shapes

    cols = 2
    rows = 4
    top = Inches(3.5)
    left = Inches(0.5)
    width = Inches(2.5)
    height = Inches(3)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2)

    table.cell(0, 0).text = 'Resume Review'
    table.cell(1, 0).text = 'Info Session 1'
    table.cell(2, 0).text = 'Info Session 2'
    table.cell(3, 0).text = 'Coffee Chat'

    for aindex,arow in attendance.iterrows():
        if NetID == str(arow['NetID']):
            table.cell(0,1).text = str(arow['Resume Review'])
            table.cell(1,1).text = str(arow['Info Session 1'])
            table.cell(2,1).text = str(arow['Info Session 2'])
            table.cell(3,1).text = str(arow['Coffee Chat'])
        else:
            table.cell(0,1).text = ''
            table.cell(1,1).text = ''
            table.cell(2,1).text = ''
            table.cell(3,1).text = ''
   
    # write column headings
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Garamond'

    
    slide = prs.slides.add_slide(blank_slide_layout)

    t2left = Inches(1)
    t2width = Inches(8)
    t2height = Inches(8)
    t2top = Inches(0.5)
    
    txBox2 = slide.shapes.add_textbox(t2left,t2top,t2width,t2height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    tf2.text = ''
    p = tf2.add_paragraph()
    p.text = 'Why GCC: ' + str(row['Why are you interested in applying for GCC']) + '\n' + '\n' + 'Definition of success: ' + str(row['What is your definition of success?'])

    p.font.size = Pt(16)
    p.font.name = 'Garamond'


os.chdir(dirname)
prs.save('test.pptx')